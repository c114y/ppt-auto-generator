import streamlit as st
from openai import OpenAI
import docx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
import concurrent.futures
import io
import re

# ================= 輔助函數 =================
def extract_text_from_docx(file):
    """提取 Word 文檔內容"""
    doc = docx.Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def clean_output(text):
    """強制清洗大模型可能生成的 Markdown 符號"""
    if text:
        return text.replace("**", "").replace("*", "")
    return ""

def set_font_style(run, font_size=18):
    """強制設置字體規範：中文黑體，數字/英文 Times New Roman"""
    run.font.name = 'Times New Roman'
    run.font.size = Pt(font_size)
    rPr = run.font._element
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'))
        rPr.append(ea)
    ea.set('typeface', '黑體')

# ================= 核心大模型調用邏輯 =================
def call_llm(prompt, client, model_name):
    """統一的大模型調用接口"""
    response = client.chat.completions.create(
        model=model_name,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2 # 極低隨機性，確保數據嚴謹、不隨意擴寫
    )
    return clean_output(response.choices[0].message.content)

def generate_outline(text, custom_req, client, model_name):
    prompt = f"""
    結合以下論文內容，幫我做一份答辯PPT大綱及每個大綱對應的內容。預計15~18頁。
    
    絕對不可違背的紅線要求（違規將導致系統崩潰）：
    1. 必須且只能使用 ====PAGE_BREAK==== 作為每頁PPT之間的分隔符（第一頁開頭不要加，放在每頁中間）。
    2. 標題直接寫具體內容，絕對不允許在標題前添加“第x頁：”或類似字樣。
    3. 必須100%忠於原文。涉及所有圖表、數據展示，必須精確抓取原文原始數據，嚴禁對數據範圍或子集進行任何形式的匯總、修改、縮減或概括。
    4. 全文輸出嚴禁出現任何英文文本或生成環境的水印提示字樣（純中文輸出）。
    5. 嚴禁在輸出中使用雙星號(**)等Markdown符號加粗。按分級標題和要點式內容生成，每點40-80字，用符號（-）分項。
    6. 客戶附加要求優先於上述規則執行：{custom_req}
    
    論文正文：
    {text[:30000]}
    """
    return call_llm(prompt, client, model_name)

def generate_speech(outline, client, model_name):
    prompt = f"""
    結合以下PPT大綱，幫我寫一份答辯演講稿。
    要求：
    1. 深度結合大綱重點，語言專業流暢，嚴禁脫離大綱自由發揮。
    2. 篇幅800字左右。
    3. 必須以“【第X頁演講詞】”作為每一段的開頭，與大綱內容一一對應。
    4. 嚴禁使用雙星號(**)等Markdown符號排版，全文不得包含任何額外的英文說明或水印詞彙。
    
    PPT大綱：
    {outline}
    """
    return call_llm(prompt, client, model_name)

def build_ppt_file(outline, template_path="template.pptx"):
    """基於模板自動化排版並生成 PPTX 文件流"""
    # 讀取你事先準備好的母版文件
    try:
        prs = Presentation(template_path)
    except Exception as e:
        raise FileNotFoundError("找不到 template.pptx 文件，請確認已將模板上傳至代碼同一目錄下！")
    
    # 使用強制定界符進行精準切分，徹底解決頁數錯亂
    pages = outline.split("====PAGE_BREAK====")
    
    for page in pages:
        if not page.strip(): continue
        
        # 調用模板中的「標題與內容」版式（索引通常為1）
        slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout) 
        
        lines = page.strip().split('\n')
        lines = [line.strip() for line in lines if line.strip()]
        if not lines: continue
        
        # 清洗標題：正則剔除可能殘留的「第X頁」等前綴
        raw_title = lines[0]
        clean_title = re.sub(r"^(第\d+頁[：:]\s*)", "", raw_title).replace("**", "")
        content_lines = lines[1:]
        
        # 寫入標題並強制設定字體
        if slide.shapes.title:
            slide.shapes.title.text = clean_title
            for par in slide.shapes.title.text_frame.paragraphs:
                for run in par.runs:
                    set_font_style(run, 28)
                
        # 智能解析正文要點並寫入模板佔位符
        if len(slide.placeholders) > 1:
            # 找到正文佔位符
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.clear() # 清除模板自帶的提示佔位符
            
            for line in content_lines:
                line = line.strip()
                if not line: continue
                
                p = tf.add_paragraph()
                # 剔除橫槓並轉換為原生層級縮進
                clean_line = line.replace("-", "").strip()
                p.text = clean_line
                p.level = 1 if line.startswith('-') else 0 
                
                for run in p.runs:
                    set_font_style(run, 18)
                
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# ================= Streamlit UI 界面 =================
st.set_page_config(page_title="自動化答辯PPT系統", layout="wide")

# 初始化會話緩存
if 'speech_text' not in st.session_state:
    st.session_state.speech_text = None
if 'ppt_io' not in st.session_state:
    st.session_state.ppt_io = None

# --- 側邊欄 API 配置 ---
st.sidebar.header("⚙️ API 配置")
api_key = st.sidebar.text_input("API Key", type="password")
base_url = st.sidebar.text_input("接口地址 (Base URL)", value="https://api.packyapi.com/v1") 
model_name = st.sidebar.text_input("模型名稱 (Model)", value="gemini-3.1-pro-preview")

st.title("⚡ 學術答辯 PPT & 演講稿自動化生成系統")

uploaded_file = st.file_uploader("📂 請上傳論文正文 (僅支持 .docx 格式)", type=["docx"])
custom_prompt = st.text_area("✍️ 客戶特殊附加要求 (選填，直接粘貼)", height=100)

if st.button("🚀 開始極速生成 (多線程)"):
    if not api_key:
        st.error("❌ 請先在左側側邊欄填寫 API Key！")
    elif not uploaded_file:
        st.warning("⚠️ 請先上傳一篇 Word 格式的論文！")
    else:
        try:
            client = OpenAI(api_key=api_key, base_url=base_url)
            
            with st.spinner("正在提取論文數據..."):
                doc_text = extract_text_from_docx(uploaded_file)
            
            st.info("第一步：正在精準提取源文檔並生成 PPT 大綱 (已啟用強制定界符與數據保真模塊)...")
            outline = generate_outline(doc_text, custom_prompt, client, model_name)
            
            with st.expander("👀 預覽生成的大綱 (點擊展開)"):
                st.write(outline)
            
            st.info("第二步：正在多線程同步生成演講稿與掛載模板排版...")
            
            with concurrent.futures.ThreadPoolExecutor() as executor:
                future_speech = executor.submit(generate_speech, outline, client, model_name)
                # 直接使用同目錄下的 template.pptx 進行掛載
                future_ppt = executor.submit(build_ppt_file, outline, "template.pptx")
                
                st.session_state.speech_text = future_speech.result()
                st.session_state.ppt_io = future_ppt.result()
                
            st.success("🎉 全部生成成功！請點擊下方按鈕下載。")
                
        except Exception as e:
            st.error(f"生成過程中出現錯誤，請檢查配置或文檔內容: {e}")

if st.session_state.speech_text and st.session_state.ppt_io:
    st.markdown("---")
    col1, col2 = st.columns(2)
    with col1:
        st.download_button(
            label="📥 下載配套演講稿 (.txt)", 
            data=st.session_state.speech_text, 
            file_name="答辯演講稿.txt", 
            mime="text/plain"
        )
    with col2:
        st.download_button(
            label="📥 下載排版後PPT (.pptx)", 
            data=st.session_state.ppt_io, 
            file_name="自動化排版答辯.pptx", 
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )